import streamlit as st
import sqlite3
import os
import datetime
import time
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import io

# ================= CONFIG =================
st.set_page_config(page_title="STAFF6 COMMAND CENTER", layout="wide")

UPLOAD_DIR = "uploads"
DB_DIR = "database"

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(DB_DIR, exist_ok=True)

# ================= LOGIN =================
ADMIN_USER = "admin06"
ADMIN_PASS = "St006904#"

if "login" not in st.session_state:
    st.session_state["login"] = False

# ================= UNITS =================
UNITS = ["พล.1 รอ.", "พล.ร.2 รอ.", "พล.ม.2 รอ.", "กรม ทย.รอ.อย."]

# ================= STATUS =================
STATUS = [
    "ยังไม่ดำเนินการ 🔴",
    "กำลังดำเนินการ 🟡",
    "เสร็จสิ้น 🟢"
]

def norm(s):
    if not s:
        return "ยังไม่ดำเนินการ 🔴"
    s = str(s)
    if "ยังไม่ดำเนิน" in s:
        return "ยังไม่ดำเนินการ 🔴"
    if "เสร็จ" in s:
        return "เสร็จสิ้น 🟢"
    if "กำลังดำเนิน" in s:
        return "กำลังดำเนินการ 🟡"
    return "ยังไม่ดำเนินการ 🔴"

# ================= DB =================
def safe(u):
    return u.replace(" ", "_").replace(".", "")

def db_path(unit):
    folder = os.path.join(DB_DIR, safe(unit))
    os.makedirs(folder, exist_ok=True)
    return os.path.join(folder, "data.db")

def connect(unit):
    conn = sqlite3.connect(db_path(unit), check_same_thread=False)
    c = conn.cursor()

    c.execute("""
    CREATE TABLE IF NOT EXISTS reports (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        unit TEXT,
        task TEXT,
        detail TEXT,
        progress INTEGER,
        status TEXT,
        problem TEXT,
        images TEXT,
        report_date TEXT,
        time TEXT
    )
    """)

    c.execute("""
    CREATE TABLE IF NOT EXISTS history (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        unit TEXT,
        task TEXT,
        detail TEXT,
        progress INTEGER,
        status TEXT,
        problem TEXT,
        images TEXT,
        report_date TEXT,
        time TEXT
    )
    """)

    conn.commit()
    return conn, c

# ================= EXPORT PPT =================
def export_ppt(data):

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ===== KPI =====
    status_count = {
        "ยังไม่ดำเนินการ 🔴": 0,
        "กำลังดำเนินการ 🟡": 0,
        "เสร็จสิ้น 🟢": 0
    }

    for d in data:
        status_count[norm(d[5])] += 1

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "STAFF6 SUMMARY"

    slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(6), Inches(3)
    ).text = f"""
TOTAL: {len(data)}
🔴 {status_count['ยังไม่ดำเนินการ 🔴']}
🟡 {status_count['กำลังดำเนินการ 🟡']}
🟢 {status_count['เสร็จสิ้น 🟢']}
"""

    # ===== DETAIL + IMAGE =====
    for d in data:

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{d[1]} | {d[2]}"

        text = f"""
หน่วย: {d[1]}
งาน: {d[2]}
รายละเอียด: {d[3]}
ความคืบหน้า: {d[4]}%
สถานะ: {norm(d[5])}
ปัญหา: {d[6]}
วันที่: {d[8]}
"""

        slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(6), Inches(4)
        ).text = text

        # 🔥 ใส่ภาพ (สำคัญ)
        if d[7]:
            imgs = d[7].split(",")

            x = 6
            y = 1

            for img in imgs:
                if img and os.path.exists(img):
                    try:
                        slide.shapes.add_picture(
                            img,
                            Inches(x),
                            Inches(y),
                            width=Inches(3)
                        )
                        x += 3
                        if x > 9:
                            x = 6
                            y += 2
                    except:
                        pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ================= USER =================
def user_app():

    st.title("📌 พื้นที่สำหรับหน่วยรายงาน")

    unit = st.selectbox("เลือกหน่วย", UNITS)
    report_date = st.date_input("📅 วันที่รายงาน", datetime.date.today())

    conn, c = connect(unit)

    task = st.text_input("งาน").strip()
    detail = st.text_area("รายละเอียด")
    progress = st.number_input("ความคืบหน้า (%)", 0, 100)
    status = st.selectbox("สถานะ", STATUS)
    problem = st.text_area("ปัญหา")

    files = st.file_uploader("📸 แนบรูป", accept_multiple_files=True)

    images = []
    if files:
        for f in files:
            name = f"{time.time()}_{f.name}"
            path = os.path.join(UPLOAD_DIR, name)
            with open(path, "wb") as w:
                w.write(f.getbuffer())
            images.append(path)

    if st.button("📤 ส่งรายงาน"):

        # history
        c.execute("""
        INSERT INTO history VALUES (NULL,?,?,?,?,?,?,?,?,?)
        """, (
            unit, task, detail, progress,
            norm(status), problem,
            ",".join(images),
            str(report_date),
            str(datetime.datetime.now())
        ))

        # latest
        existing = c.execute("""
        SELECT id, progress, images FROM reports
        WHERE unit=? AND task=?
        ORDER BY id DESC LIMIT 1
        """, (unit, task)).fetchone()

        if existing:
            rid, old_progress, old_images = existing

            new_progress = max(old_progress, progress)

            old_list = old_images.split(",") if old_images else []
            new_list = old_list + images

            c.execute("""
            UPDATE reports
            SET detail=?, progress=?, status=?, problem=?, images=?, report_date=?, time=?
            WHERE id=?
            """, (
                detail,
                new_progress,
                norm(status),
                problem,
                ",".join(new_list),
                str(report_date),
                str(datetime.datetime.now()),
                rid
            ))
        else:
            c.execute("""
            INSERT INTO reports VALUES (NULL,?,?,?,?,?,?,?,?,?)
            """, (
                unit, task, detail, progress,
                norm(status), problem,
                ",".join(images),
                str(report_date),
                str(datetime.datetime.now())
            ))

        conn.commit()
        st.success("บันทึกเรียบร้อย")

    st.stop()

# ================= LOAD =================
def load_history():
    data = []
    for u in UNITS:
        conn, c = connect(u)
        rows = c.execute("SELECT * FROM history").fetchall()
        for r in rows:
            r = list(r)
            r[5] = norm(r[5])
            data.append(r)
    return data

def load_latest():
    data = []
    for u in UNITS:
        conn, c = connect(u)
        rows = c.execute("SELECT * FROM reports").fetchall()
        for r in rows:
            r = list(r)
            r[5] = norm(r[5])
            data.append(r)
    return data

# ================= ADMIN =================
def admin_app():

    st.title("🚨 STAFF6 COMMAND CENTER")

    history = load_history()
    latest = load_latest()

    # FIX DATE RANGE
    all_dates = []
    for d in history:
        try:
            all_dates.append(datetime.datetime.strptime(d[8], "%Y-%m-%d").date())
        except:
            pass

    min_date = min(all_dates) if all_dates else datetime.date.today()
    max_date = max(all_dates) if all_dates else datetime.date.today()

    with st.sidebar:
        if st.button("🚪 Logout"):
            st.session_state["login"] = False
            st.rerun()

        unit_filter = st.selectbox("หน่วย", ["ทั้งหมด"] + UNITS)
        from_date = st.date_input("From", min_date)
        to_date = st.date_input("To", max_date)

    # FILTER
    filtered_history = []
    for d in history:
        try:
            dd = datetime.datetime.strptime(d[8], "%Y-%m-%d").date()
        except:
            continue

        if unit_filter != "ทั้งหมด" and d[1] != unit_filter:
            continue

        if not (from_date <= dd <= to_date):
            continue

        filtered_history.append(d)

    filtered_latest = []
    for d in latest:
        if unit_filter != "ทั้งหมด" and d[1] != unit_filter:
            continue
        filtered_latest.append(d)

    # KPI
    st.subheader("📊 KPI")

    status_list = [norm(x[5]) for x in filtered_latest]

    c1,c2,c3,c4 = st.columns(4)
    c1.metric("📦 ทั้งหมด", len(filtered_latest))
    c2.metric("🟡 กำลังดำเนินการ", status_list.count("กำลังดำเนินการ 🟡"))
    c3.metric("🟢 เสร็จสิ้น", status_list.count("เสร็จสิ้น 🟢"))
    c4.metric("🔴 ยังไม่ดำเนินการ", status_list.count("ยังไม่ดำเนินการ 🔴"))

    st.markdown("---")

    # GRAPH
    st.subheader("📈 ความคืบหน้า")

    if filtered_latest:
        df = pd.DataFrame(filtered_latest, columns=[
            "ID","หน่วย","งาน","รายละเอียด","%","สถานะ",
            "ปัญหา","รูป","วันที่","เวลา"
        ])

        if unit_filter == "ทั้งหมด":
            st.bar_chart(df.groupby("หน่วย")["%"].mean())
        else:
            st.bar_chart(df.groupby("งาน")["%"].mean())

    st.markdown("---")

    # REPORT
    st.subheader("📄 รายงาน")

    for i, d in enumerate(filtered_history):

        col1,col2 = st.columns([3,1])

        with col1:
            st.markdown(f"""
### 🏷 {d[1]} | {d[2]} | {norm(d[5])}

{d[3]}

📊 {d[4]}%  
⚠️ {d[6] if d[6] else "-"}  
📅 {d[8]}  
🕒 {d[9]}
""")

            if d[7]:
                for img in d[7].split(","):
                    if img and os.path.exists(img):
                        st.image(img, width=250)

        with col2:
            if st.button("🗑 ลบ", key=f"{i}_{d[0]}"):
                conn, c = connect(d[1])
                c.execute("DELETE FROM reports WHERE id=?", (d[0],))
                conn.commit()
                st.rerun()

    st.markdown("---")

    if st.button("📤 Export PPT"):
        ppt = export_ppt(filtered_history)
        st.download_button("📥 ดาวน์โหลด", ppt, file_name="STAFF6_REPORT.pptx")

# ================= LOGIN =================
def login_page():
    st.title("🔐 STAFF6 LOGIN")

    u = st.text_input("User")
    p = st.text_input("Password", type="password")

    if st.button("Login"):
        if u == ADMIN_USER and p == ADMIN_PASS:
            st.session_state["login"] = True
            st.rerun()
        else:
            st.error("Login ไม่ถูกต้อง")

# ================= MAIN =================
def main():
    if st.session_state["login"]:
        admin_app()
    else:
        login_page()
        user_app()

main()
